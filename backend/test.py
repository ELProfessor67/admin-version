from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE

def shrink_and_wrap_ppt(input_path, output_path, scale=0.6):
    prs = Presentation(input_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue

            text_frame = shape.text_frame
            text_frame.word_wrap = True  # like white-space: normal
            text_frame.auto_size = MSO_AUTO_SIZE.NONE  # disable autofit safely

            for para in text_frame.paragraphs:
                for run in para.runs:
                    if run.font.size:
                        # reduce font size by 40%
                        run.font.size = Pt(run.font.size.pt * scale)

    prs.save(output_path)
    print(f"✅ Saved fixed file: {output_path}")

# Example usage
shrink_and_wrap_ppt("example.pptx", "example_fixed.pptx")
